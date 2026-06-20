"""LSTM vs Transformer 발표자료 v2 — tiny-budget-ada 디자인 문법 (python-pptx).
본문: 한 텍스트 박스에 [헤더(blue) + PowerPoint 기본 불릿 항목]. 한 불릿 안에서 부분 강조(단어/구) 가능.
  PowerPoint에서 Tab 들여쓰기·자동 불릿 동작 → 직접 편집 쉬움.
실행: /Users/jihunjang/miniconda3/bin/python redesign/build.py
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

C={'ink':'111827','body':'374151','sub':'6B7280','blue':'1D4ED8','green':'0F766E','line':'E5E7EB','mut':'BFC6D0'}
def rgb(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
EMU=914400
MX=0.5; W=13.333-MX*2
FIG="/Users/jihunjang/workspace/ust/deeplearning2/term-project-lstm-vs-transformer-agnews/redesign"
OUT="/Users/jihunjang/workspace/ust/deeplearning2/term-project-lstm-vs-transformer-agnews/deliverables/5_presentation_v2.pptx"
STEPS=["배경","방법","결과","결론"]; TOT=10

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def tbox(s,x,y,w,h):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=MSO_ANCHOR.TOP
    return tf
def run(p,txt,b=False,c='374151',sz=14,face='Pretendard'):
    r=p.add_run(); r.text=txt; r.font.bold=b; r.font.size=Pt(sz); r.font.name=face; r.font.color.rgb=rgb(c); return r
def space(p,before=None,after=None,line=None):
    pPr=p._p.get_or_add_pPr()
    if line is not None:
        e=etree.SubElement(pPr,qn('a:lnSpc')); x=etree.SubElement(e,qn('a:spcPct')); x.set('val',str(int(line*100000)))
    if before is not None:
        e=etree.SubElement(pPr,qn('a:spcBef')); x=etree.SubElement(e,qn('a:spcPts')); x.set('val',str(int(before*100)))
    if after is not None:
        e=etree.SubElement(pPr,qn('a:spcAft')); x=etree.SubElement(e,qn('a:spcPts')); x.set('val',str(int(after*100)))
def bullet(p,lvl=0,char='●',clr='1D4ED8'):
    pPr=p._p.get_or_add_pPr()
    pPr.set('marL',str(int(0.34*(lvl+1)*EMU))); pPr.set('indent',str(-int(0.26*EMU)))
    e=etree.SubElement(pPr,qn('a:buClr')); sc=etree.SubElement(e,qn('a:srgbClr')); sc.set('val',clr)
    e=etree.SubElement(pPr,qn('a:buSzPct')); e.set('val','76000')
    e=etree.SubElement(pPr,qn('a:buFont')); e.set('typeface','Arial')
    e=etree.SubElement(pPr,qn('a:buChar')); e.set('char',char)
def hline(s,x,y,w,clr='E5E7EB',wpt=1.0):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x),Inches(y),Inches(x+w),Inches(y))
    ln.line.color.rgb=rgb(clr); ln.line.width=Pt(wpt); return ln
def rect(s,x,y,w,h,clr):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(clr); sh.line.fill.background(); sh.shadow.inherit=False; return sh
def pic(s,path,x,y,bw,bh,dpi=200):
    iw,ih=Image.open(path).size; iwin=iw/dpi; ihin=ih/dpi
    sc=min(bw/iwin,bh/ihin); w=iwin*sc; h=ihin*sc
    s.shapes.add_picture(path,Inches(x+(bw-w)/2),Inches(y+(bh-h)/2),Inches(w),Inches(h))
def page(s,n):
    tf=tbox(s,11.9,7.12,1.1,0.26); p=tf.paragraphs[0]; run(p,f"{n} / {TOT}",False,'AAB2BD',9.5); p.alignment=PP_ALIGN.RIGHT
def head(s,active,headruns):
    tf=tbox(s,MX,0.36,9,0.3); p=tf.paragraphs[0]
    for i,st in enumerate(STEPS):
        run(p,st,i==active,'1D4ED8' if i==active else 'BFC6D0',12.5)
        if i<len(STEPS)-1: run(p,'     ',False,'BFC6D0',12.5)
    tf2=tbox(s,MX,0.74,W,0.55); p2=tf2.paragraphs[0]; space(p2,line=1.0)
    for txt,c in headruns: run(p2,txt,True,c,21)
def body(s,x,y,w,h,blocks,fs=14,gap=14):
    tf=tbox(s,x,y,w,h); first=[True]
    def para():
        if first[0]: first[0]=False; return tf.paragraphs[0]
        return tf.add_paragraph()
    for bi,blk in enumerate(blocks):
        if blk.get('h'):
            p=para(); space(p,before=(gap if bi>0 else 0),after=7); run(p,blk['h'],True,'1D4ED8',fs+1.5)
        for it in blk['items']:
            lv=it.get('lv',0) if isinstance(it,dict) else 0
            parts=it['parts'] if isinstance(it,dict) else it
            p=para(); space(p,after=5,line=1.25); bullet(p,lvl=lv,char=('–' if lv else '●'),clr=('6B7280' if lv else '1D4ED8'))
            for pt in parts: run(p,pt[0],pt[1] if len(pt)>1 else False,pt[2] if len(pt)>2 else '374151',fs)

# ── 1 표지 (메타 프레임)
s=slide(); rect(s,0,0,0.16,7.5,'1D4ED8')
tf=tbox(s,0.9,0.68,8,0.35); run(tf.paragraphs[0],"DEEP LEARNING · TERM PROJECT",True,'1D4ED8',14)
tf=tbox(s,9.4,0.68,3,0.35); p=tf.paragraphs[0]; run(p,"2026",False,'6B7280',14); p.alignment=PP_ALIGN.RIGHT
hline(s,0.9,1.2,11.5)
tf=tbox(s,0.9,2.0,11.5,1.0); p=tf.paragraphs[0]; run(p,"LSTM ",True,'111827',50); run(p,"vs",True,'1D4ED8',50); run(p," Transformer",True,'111827',50)
tf=tbox(s,0.9,3.18,11.5,0.7); run(tf.paragraphs[0],"Encoder, 텍스트 분류에서의 통제된 비교",True,'111827',31)
tf=tbox(s,0.9,4.22,11.5,0.5); run(tf.paragraphs[0],"AG News 4-class 주제 분류, 두 시퀀스 인코더를 동일 조건에서 비교",False,'6B7280',19)
hline(s,0.9,5.42,11.5)
tf=tbox(s,0.9,5.74,1.55,0.35); run(tf.paragraphs[0],"발표자",True,'1D4ED8',15)
tf=tbox(s,2.45,5.72,9,0.4); p=tf.paragraphs[0]; run(p,"장지훈",True,'111827',17); run(p,"   (ETRI School · 02521122)",False,'6B7280',13)
tf=tbox(s,0.9,6.38,1.55,0.35); run(tf.paragraphs[0],"소속",True,'1D4ED8',15)
tf=tbox(s,2.45,6.36,9,0.35); run(tf.paragraphs[0],"과학기술연합대학원대학교 ETRI 스쿨",False,'374151',16,'Pretendard SemiBold')

# ── 2 INTRODUCTION
s=slide(); head(s,0,[("정확도 경쟁이 아니라, ",'111827'),("두 구조가 왜 다른지",'1D4ED8'),("를 규명",'111827')])
body(s,MX,1.55,W,5.2,[
  {'h':'과제','items':[
    [("AG News 4-class 뉴스 주제 분류",True,'111827'),(" (World·Sports·Business·Sci-Tech)",)],
    [("LSTM과 Transformer Encoder를 모두 from scratch로 학습",)]]},
  {'h':'방법','items':[
    [("데이터·전처리·어휘·시퀀스 길이·학습 예산",True,'111827'),("을 동일하게 고정",)],
    [("encoder 구조만 유일한 변수로 둔다",)]]},
  {'h':'질문','items':[
    [("성능·수렴·data efficiency·오분류 양상이 어떻게 다른가",)],
    [("그리고 그 차이의 원인은 무엇인가",)]]},
],15,15); page(s,2)

# ── 3 DATASET
s=slide(); head(s,0,[("balanced 4-class, 어휘는 ",'111827'),("train에서만 구축",'1D4ED8'),(" (leakage 방지)",'111827')])
body(s,MX,1.55,W,5.2,[
  {'h':'분할 (seed 42)','items':[
    [("Train 108,000 / Validation 12,000 / Test 7,600",True,'111827')]]},
  {'h':'클래스','items':[
    [("World · Sports · Business · Sci-Tech",True,'111827'),(" 4개, balanced",)],
    [("train 약 27,000 / test 1,900 each",)]]},
  {'h':'전처리 (두 모델 공통)','items':[
    [("단어 단위 tokenization (소문자화)",)],
    [("vocabulary: ",),("train-only",True,'111827'),(", 상한 20,000, min frequency 2",)],
    [("max length 128, padding 위치는 masking으로 무시",)]]},
],15,14); page(s,3)

# ── 4 MODELS
s=slide(); head(s,1,[("공유 골격, ",'111827'),("encoder만 교체",'1D4ED8')])
body(s,MX,1.55,W,5.2,[
  {'h':'공유 골격','items':[
    [("Embedding(128) → Encoder → masked mean pooling → Linear(4)",True,'111827')],
    [("encoder 한 조각만 교체, 나머지는 모두 동일",)]]},
  {'h':'encoder 비교','items':[
    [("LSTM",True,'111827'),("  bidirectional · 2 layers · hidden 128 / output 256 / params 3,220,484",)],
    [("Transformer",True,'111827'),("  2 layers · 4 heads · FFN 256 · sinusoidal PE / dim 128 / params 2,825,476",)]]},
  {'h':'파라미터 규모','items':[
    [("임베딩이 지배적이라 두 모델 파라미터는 ",),("약 14% 차",True,'111827'),("로 비슷",)]]},
],15,15); page(s,4)

# ── 5 EXPERIMENTS
s=slide(); head(s,1,[("유일한 변수 = ",'111827'),("encoder 구조",'1D4ED8'),(" (나머지 전부 고정)",'111827')])
body(s,MX,1.55,W,5.2,[
  {'h':'고정 (fair comparison)','items':[
    [("split · tokenizer · vocabulary · max length · pooling",)],
    [("optimizer · learning rate · batch size · epoch · seed",)]]},
  {'h':'학습','items':[
    [("Adam, learning rate 0.001, batch 64, 최대 8 epochs, dropout 0.1, seed 42",)],
    [("model selection은 ",),("validation",True,'111827'),(", test는 1회",)]]},
  {'h':'Ablation · 지표','items':[
    [("training set 5 / 25 / 50 / 100%",True,'111827'),(" (vocabulary 고정, stratified, val·test 전체)",)],
    [("accuracy · macro F1-score · loss curve · confusion matrix",)]]},
],15,14); page(s,5)

# ── 6 RESULTS loss
s=slide(); head(s,2,[("Transformer ",'111827'),("0.910",'1D4ED8'),(" vs LSTM ",'111827'),("0.833",'6B7280'),(", 격차의 원인은 overfitting",'111827')])
pic(s,FIG+"/fig_loss.png",MX,1.5,W,4.25)
body(s,MX,5.95,W,1.4,[
  {'items':[
    [("Transformer",True,'1D4ED8'),("  test accuracy ",),("0.910",True,'1D4ED8'),(" · macro F1 0.909",)],
    [("LSTM",True,'6B7280'),("  0.833 · macro F1 0.835, ",),("val loss 폭증",True,'C0392B'),("으로 best epoch 2 저장",)]]},
],14,0); page(s,6)

# ── 7 RESULTS confusion
s=slide(); head(s,2,[("두 모델 공통 ",'111827'),("Business ↔ Sci/Tech 혼동",'1D4ED8')])
pic(s,FIG+"/fig_confusion.png",MX,1.5,W,4.0)
body(s,MX,5.75,W,1.5,[
  {'items':[
    [("공통 오답 428건",True,'111827'),("  61%가 Business/Sci-Tech 경계 (의미 중첩에 따른 모호성)",)],
    [("LSTM 단독 841건",True,'111827'),("  전 class에 분산, Transformer 단독(259건)의 ",),("3.2배",True,'111827')]]},
],14,0); page(s,7)

# ── 8 ABLATION (좌 그림 + 우 body)
s=slide(); head(s,2,[("LSTM은 ",'111827'),("25%에서 포화",'6B7280'),(", Transformer는 ",'111827'),("계속 향상",'1D4ED8')])
pic(s,FIG+"/fig_ablation.png",MX,1.7,7.1,4.8)
body(s,7.9,1.75,W-7.4,4.8,[
  {'h':'Transformer','items':[
    [("데이터에 따라 단조 향상",True,'1D4ED8')],
    [("scaling: 더 줄수록 더 좋아짐",)]]},
  {'h':'LSTM','items':[
    [("25%에서 정점 후 정체",True,'111827')],
    [("overfitting으로 추가 데이터를 못 씀",)]]},
  {'h':'가설 기각','items':[
    [("'data-hungry' 가설 기각",)],
    [("전 구간에서 Transformer 우위",)]]},
],14,16); page(s,8)

# ── 9 EXTENSION mechanism
s=slide(); head(s,2,[("bag-of-words 과제, 격차는 ",'111827'),("순서가 아닌 robustness",'1D4ED8')])
pic(s,FIG+"/fig_mechanism.png",MX,1.45,W,3.6)
body(s,MX,5.25,W,1.6,[
  {'items':[
    [("PE on/off",True,'111827'),("  with-PE 0.910 vs no-PE 0.911 (차이 -0.001) → 순서 정보가 거의 불필요",)],
    [("robustness",True,'111827'),("  LSTM은 비주제어 'giddy'에 쏠려 오답, Transformer는 주제 토큰에 분산해 정답",)]]},
],14,0); page(s,9)

# ── 10 CONCLUSION
s=slide(); head(s,3,[("동일 조건에서 ",'111827'),("Transformer 우수",'1D4ED8'),(", 원인은 generalization",'111827')])
body(s,MX,1.55,W,5.2,[
  {'h':'성능','items':[
    [("동일 조건에서 test accuracy ",),("0.910 vs 0.833",True,'111827'),(" (encoder 구조만 바꾼 통제된 비교)",)]]},
  {'h':'원인','items':[
    [("generalization 차이",True,'111827'),(": LSTM은 overfitting하여 25%에서 포화, Transformer는 안정적으로 확장",)]]},
  {'h':'과제 성격','items':[
    [("bag-of-words",True,'111827'),(": positional encoding 제거해도 정확도 동일 → 격차는 순서가 아닌 robustness",)]]},
  {'h':'한계 · 향후','items':[
    [("단일 seed → 다중 seed 재확인, LSTM regularization 강화",)],
    [("순서가 중요한 과제에서 positional encoding 재검증",)]]},
],15,14); page(s,10)

prs.save(OUT); print("saved", OUT)
