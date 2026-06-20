"""LSTM vs Transformer presentation — ENGLISH version (tiny-budget-ada design grammar, python-pptx).
Mirror of build.py with English text + English figures (fig_*_en.png).
Run: /Users/jihunjang/miniconda3/bin/python redesign/build_en.py
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

C={'ink':'111827','body':'374151','sub':'6B7280','blue':'1D4ED8','green':'0F766E','line':'E5E7EB','mut':'BFC6D0'}
def rgb(h): return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
EMU=914400
MX=0.5; W=13.333-MX*2
FIG="/Users/jihunjang/workspace/ust/deeplearning2/term-project-lstm-vs-transformer-agnews/redesign"
OUT="/Users/jihunjang/workspace/ust/deeplearning2/term-project-lstm-vs-transformer-agnews/deliverables_en/4_final_presentation.pptx"
STEPS=["Background","Method","Results","Conclusion"]; TOT=10

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def tbox(s,x,y,w,h):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=MSO_ANCHOR.TOP
    return tf
def run(p,txt,b=False,c='374151',sz=14,face='Pretendard'):
    r=p.add_run(); r.text=txt; r.font.bold=b; r.font.size=Pt(sz); r.font.name=face; r.font.color.rgb=rgb(c); return r
def space(p,before=None,after=None,line=None):
    pPr=p._p.get_or_add_pPr()
    if line is not None:
        e=etree.SubElement(pPr,qn('a:lnSpc')); x=etree.SubElement(e,qn('a:spcPct')); x.set('val',str(int(line*100000)))
    if before is not None:
        e=etree.SubElement(pPr,qn('a:spcBef')); x=etree.SubElement(e,qn('a:spcPts')); x.set('val',str(int(before*100)))
    if after is not None:
        e=etree.SubElement(pPr,qn('a:spcAft')); x=etree.SubElement(e,qn('a:spcPts')); x.set('val',str(int(after*100)))
def bullet(p,lvl=0):
    pPr=p._p.get_or_add_pPr()
    pPr.set('marL',str(int(0.3*(lvl+1)*EMU))); pPr.set('indent',str(-int(0.3*EMU)))
    e=etree.SubElement(pPr,qn('a:buFont')); e.set('typeface','Arial')
    e=etree.SubElement(pPr,qn('a:buChar')); e.set('char',('•' if lvl==0 else '–'))
def hline(s,x,y,w,clr='E5E7EB',wpt=1.0):
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x),Inches(y),Inches(x+w),Inches(y))
    ln.line.color.rgb=rgb(clr); ln.line.width=Pt(wpt); return ln
def rect(s,x,y,w,h,clr):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(clr); sh.line.fill.background(); sh.shadow.inherit=False; return sh
def pic(s,path,x,y,bw,bh,dpi=200):
    iw,ih=Image.open(path).size; iwin=iw/dpi; ihin=ih/dpi
    sc=min(bw/iwin,bh/ihin); w=iwin*sc; h=ihin*sc
    s.shapes.add_picture(path,Inches(x+(bw-w)/2),Inches(y+(bh-h)/2),Inches(w),Inches(h))
def page(s,n):
    tf=tbox(s,11.9,7.12,1.1,0.26); p=tf.paragraphs[0]; run(p,f"{n} / {TOT}",False,'AAB2BD',9.5); p.alignment=PP_ALIGN.RIGHT
def head(s,active,headruns):
    tf=tbox(s,MX,0.36,11,0.3); p=tf.paragraphs[0]
    for i,st in enumerate(STEPS):
        run(p,st,i==active,'1D4ED8' if i==active else 'BFC6D0',12.5)
        if i<len(STEPS)-1: run(p,'     ',False,'BFC6D0',12.5)
    tf2=tbox(s,MX,0.74,W,0.6); p2=tf2.paragraphs[0]; space(p2,line=1.0)
    for txt,c in headruns: run(p2,txt,True,c,21)
def body(s,x,y,w,h,blocks,fs=14,gap=14):
    tf=tbox(s,x,y,w,h); first=[True]
    def para():
        if first[0]: first[0]=False; return tf.paragraphs[0]
        return tf.add_paragraph()
    for bi,blk in enumerate(blocks):
        if blk.get('h'):
            p=para(); space(p,before=(gap if bi>0 else 0),after=7); run(p,blk['h'],True,'1D4ED8',fs+1.5)
        for it in blk['items']:
            lv=it.get('lv',0) if isinstance(it,dict) else 0
            parts=it['parts'] if isinstance(it,dict) else it
            p=para(); space(p,after=5,line=1.25); bullet(p,lvl=lv)
            for pt in parts: run(p,pt[0],pt[1] if len(pt)>1 else False,pt[2] if len(pt)>2 else '374151',fs)

# ── 1 Title
s=slide(); rect(s,0,0,0.16,7.5,'1D4ED8')
tf=tbox(s,0.9,0.68,8,0.35); run(tf.paragraphs[0],"DEEP LEARNING · TERM PROJECT",True,'1D4ED8',14)
tf=tbox(s,9.4,0.68,3,0.35); p=tf.paragraphs[0]; run(p,"2026",False,'6B7280',14); p.alignment=PP_ALIGN.RIGHT
hline(s,0.9,1.2,11.5)
tf=tbox(s,0.9,2.25,11.5,1.0); p=tf.paragraphs[0]; run(p,"LSTM ",True,'111827',50); run(p,"vs",True,'1D4ED8',50); run(p," Transformer",True,'111827',50)
tf=tbox(s,0.9,3.5,11.5,0.6); run(tf.paragraphs[0],"Two sequence encoders compared on AG News 4-class topic classification",False,'6B7280',19)
hline(s,0.9,5.42,11.5)
tf=tbox(s,0.9,5.74,1.6,0.35); run(tf.paragraphs[0],"Team",True,'1D4ED8',15)
tf=tbox(s,2.5,5.72,9,0.4); p=tf.paragraphs[0]; run(p,"Jihun Jang",True,'111827',17); run(p,"   (ETRI School · 02521122)",False,'6B7280',13)
tf=tbox(s,0.9,6.38,1.6,0.35); run(tf.paragraphs[0],"Affiliation",True,'1D4ED8',15)
tf=tbox(s,2.5,6.36,9,0.35); run(tf.paragraphs[0],"University of Science and Technology · ETRI School",False,'374151',16,'Pretendard SemiBold')

# ── 2 INTRODUCTION
s=slide(); head(s,0,[("Not an accuracy race: ",'111827'),("we explain WHY the two architectures differ",'1D4ED8')])
body(s,MX,1.6,W,5.2,[
  {'h':'Task','items':[
    [("AG News 4-class news topic classification",True,'111827'),(" (World·Sports·Business·Sci-Tech)",)],
    [("Both LSTM and Transformer Encoder trained from scratch",)]]},
  {'h':'Method','items':[
    [("Data, preprocessing, vocabulary, sequence length, training budget",True,'111827'),(" held identical",)],
    [("Only the encoder differs by architecture",)]]},
  {'h':'Question','items':[
    [("How do performance, convergence, data efficiency, and error patterns differ?",)],
    [("And what causes the difference?",)]]},
],15,15); page(s,2)

# ── 3 DATASET
s=slide(); head(s,0,[("Balanced 4-class, vocabulary ",'111827'),("built from train only",'1D4ED8'),(" (prevents leakage)",'111827')])
body(s,MX,1.6,W,5.2,[
  {'h':'Split (seed 42)','items':[
    [("Train 108,000 / Validation 12,000 / Test 7,600",True,'111827')]]},
  {'h':'Classes','items':[
    [("World · Sports · Business · Sci-Tech",True,'111827'),(", balanced",)],
    [("train ~27,000 / test 1,900 each",)]]},
  {'h':'Preprocessing (shared by both models)','items':[
    [("word-level tokenization (lowercased)",)],
    [("vocabulary: ",),("train-only",True,'111827'),(", max 20,000, min frequency 2",)],
    [("max length 128, padding positions masked out",)]]},
],15,14); page(s,3)

# ── 4 MODELS
s=slide(); head(s,1,[("Shared backbone, ",'111827'),("only the encoder mechanism differs",'1D4ED8')])
body(s,MX,1.6,W,5.2,[
  {'h':'Shared backbone','items':[
    [("Embedding(128) → Encoder → masked mean pooling → Linear(4)",True,'111827')],
    [("only the encoder is swapped, everything else identical",)]]},
  {'h':'Core mechanism (basis for interpretation)','items':[
    [("LSTM",True,'6B7280'),("  processes tokens sequentially, accumulating context in the hidden state (recurrence)",)],
    [("Transformer",True,'1D4ED8'),("  connects all tokens directly via self-attention, distributes evidence across tokens",)]]},
  {'h':'Setup · scale','items':[
    [("LSTM",True,'6B7280'),("  bidirectional 2 layers · hidden 128 · params 3,220,484",)],
    [("Transformer",True,'1D4ED8'),("  2 layers · 4 heads · FFN 256 · sinusoidal PE · params 2,825,476",)],
    [("embedding dominates, so the two models are within ",),("~14%",True,'111827'),(" in parameters",)]]},
],14,13); page(s,4)

# ── 5 EXPERIMENTS
s=slide(); head(s,1,[("Only variable = ",'111827'),("encoder architecture",'1D4ED8'),(" (everything else fixed)",'111827')])
body(s,MX,1.6,W,5.2,[
  {'h':'Fixed (fair comparison)','items':[
    [("split · tokenizer · vocabulary · max length · pooling",)],
    [("optimizer · learning rate · batch size · epoch · seed",)]]},
  {'h':'Training','items':[
    [("Adam, learning rate 0.001, batch 64, up to 8 epochs, dropout 0.1, seed 42",)],
    [("model selection on ",),("validation",True,'111827'),(", test evaluated once",)]]},
  {'h':'Ablation · metrics','items':[
    [("training set 5 / 25 / 50 / 100%",True,'111827'),(" (vocabulary fixed, stratified, full val·test)",)],
    [("accuracy · macro F1-score · loss curve · confusion matrix",)]]},
],15,14); page(s,5)

# ── 6 RESULTS loss
s=slide(); head(s,2,[("Transformer ",'111827'),("0.910",'1D4ED8'),(" vs LSTM ",'111827'),("0.833",'6B7280'),(", the gap comes from overfitting",'111827')])
pic(s,FIG+"/fig_loss_en.png",MX,1.55,W,4.25)
body(s,MX,5.95,W,1.4,[
  {'items':[
    [("Transformer",True,'1D4ED8'),("  test accuracy ",),("0.910",True,'1D4ED8'),(" · macro F1 0.909",)],
    [("LSTM",True,'6B7280'),("  0.833 · macro F1 0.835, ",),("val loss explodes",True,'C0392B'),(", best epoch 2 saved",)]]},
],14,0); page(s,6)

# ── 7 RESULTS confusion
s=slide(); head(s,2,[("Both models confuse ",'111827'),("Business ↔ Sci/Tech the most",'1D4ED8')])
pic(s,FIG+"/fig_confusion_en.png",MX,1.55,W,4.0)
body(s,MX,5.8,W,1.5,[
  {'items':[
    [("Both wrong: 428",True,'111827'),("  61% on the Business / Sci-Tech boundary (semantic overlap)",)],
    [("LSTM-only: 841",True,'111827'),("  spread across all classes, ",),("3.2×",True,'111827'),(" the Transformer-only (259)",)]]},
],14,0); page(s,7)

# ── 8 ABLATION
s=slide(); head(s,2,[("LSTM ",'111827'),("saturates at 25%",'6B7280'),(", Transformer ",'111827'),("keeps improving",'1D4ED8')])
pic(s,FIG+"/fig_ablation_en.png",MX,1.7,7.1,4.8)
body(s,7.9,1.8,W-7.4,4.8,[
  {'h':'Transformer','items':[
    [("monotonic increase with data (scaling)",True,'1D4ED8')],
    [("more data, better accuracy",)]]},
  {'h':'LSTM','items':[
    [("peaks at 25%, then saturation",True,'6B7280')],
    [("overfitting prevents using more data",)]]},
  {'h':'Hypothesis rejected','items':[
    [("'data-hungry' hypothesis rejected",)],
    [("Transformer leads across all data sizes",)]]},
],14,16); page(s,8)

# ── 9 EXTENSION mechanism
s=slide(); head(s,2,[("Bag-of-words task: the gap is ",'111827'),("robustness, not word order",'1D4ED8')])
tf=tbox(s,MX,1.4,W,0.8); p=tf.paragraphs[0]; space(p,line=1.18)
run(p,"Input (label: ",False,'6B7280',12.5); run(p,"Sports",True,'1D4ED8',12.5); run(p,")   ",False,'6B7280',12.5)
run(p,'"Giddy Phelps Touches Gold for First Time Michael Phelps won the gold medal in the 400 individual medley and set a world record in a time of 4 minutes 8.26 seconds."',False,'374151',12.5)
pic(s,FIG+"/fig_mechanism_en.png",MX,2.28,W,2.95)
body(s,MX,5.42,W,1.55,[
  {'items':[
    [("PE on/off",True,'111827'),("  with-PE 0.910 vs no-PE 0.911 (diff -0.001), word order barely matters",)],
    [("robustness",True,'111827'),("  LSTM relies on off-topic 'giddy' and misclassifies; Transformer distributes evidence over topical tokens and classifies correctly",)]]},
],14,0); page(s,9)

# ── 10 CONCLUSION
s=slide(); head(s,3,[("Under identical conditions, ",'111827'),("Transformer wins",'1D4ED8'),(", cause is generalization",'111827')])
body(s,MX,1.6,W,5.2,[
  {'h':'Performance','items':[
    [("test accuracy ",),("0.910 vs 0.833",True,'111827'),(" under identical conditions (controlled comparison, only the encoder changed)",)]]},
  {'h':'Cause (mechanism difference)','items':[
    [("Transformer",True,'1D4ED8'),(": self-attention distributes evidence across tokens → robust, stable scaling",)],
    [("LSTM",True,'6B7280'),(": recurrence relies on specific tokens → overfitting, saturation at 25%",)],
    [("the generalization gap between the two mechanisms drives the performance gap",)]]},
  {'h':'Task nature','items':[
    [("bag-of-words",True,'111827'),(": removing positional encoding leaves accuracy unchanged, the gap is robustness not order",)]]},
  {'h':'Limitations · future','items':[
    [("single seed → confirm with multiple seeds, strengthen LSTM regularization",)],
    [("re-examine positional encoding on order-sensitive tasks",)]]},
],14,12); page(s,10)

prs.save(OUT); print("saved", OUT)
